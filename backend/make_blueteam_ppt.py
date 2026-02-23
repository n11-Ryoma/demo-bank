from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path("BlueTeam_Exercise_Environment_Overview.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLOR_BG = RGBColor(246, 250, 255)
COLOR_ACCENT = RGBColor(0, 92, 169)
COLOR_TEXT = RGBColor(30, 30, 30)
COLOR_SUB = RGBColor(90, 90, 90)
COLOR_BOX = RGBColor(255, 255, 255)
COLOR_BOX_LINE = RGBColor(190, 208, 226)
COLOR_WEB = RGBColor(230, 245, 255)
COLOR_APP = RGBColor(235, 250, 238)
COLOR_DB = RGBColor(255, 245, 230)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.32))
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_ACCENT
    band.line.fill.background()


def add_header(slide, title, subtitle=""):
    add_bg(slide)
    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12.2), Inches(0.9))
    tf = t.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.name = "Meiryo"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(12.0), Inches(0.5))
        stf = s.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.name = "Meiryo"
        sp.font.size = Pt(18)
        sp.font.color.rgb = COLOR_SUB


def add_bullets(slide, x, y, w, h, lines, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Meiryo"
        p.font.size = Pt(size)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(8)


def add_card(slide, x, y, w, h, title, lines, fill_color=COLOR_BOX):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = COLOR_BOX_LINE

    title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.35), Inches(0.5))
    ttf = title_box.text_frame
    ttf.text = title
    tp = ttf.paragraphs[0]
    tp.font.name = "Meiryo"
    tp.font.size = Pt(18)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_ACCENT

    add_bullets(slide, x + 0.2, y + 0.58, w - 0.35, h - 0.7, lines, size=15)


def add_line(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLOR_ACCENT
    line.line.width = Pt(2)


def add_label(slide, x, y, text, size=13):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.8), Inches(0.35))
    tf = b.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name = "Meiryo"
    p.font.size = Pt(size)
    p.font.color.rgb = COLOR_SUB


# 1. Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "環境説明資料", "demo-bank backend (Apache / Spring Boot / PostgreSQL)")
add_bullets(slide, 0.9, 2.15, 12.0, 3.8, [
    "この資料は『どういう環境か』『どこにログが出るか』を共有するためのものです。",
    "対象: Webサーバ、業務ロジック、DBの3層。",
    "用途: 監視設定、障害切り分け、運用引き継ぎ。",
], size=22)


# 2. Environment diagram
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "1. 環境構成図", "3層構成と主要ログ保管場所")

add_card(
    slide, 0.7, 2.0, 3.8, 2.2,
    "Web: Apache",
    [
        "公開エンドポイントの入口",
        "access/errorを記録",
        "Path: /var/log/apache2",
    ],
    fill_color=COLOR_WEB,
)

add_card(
    slide, 4.9, 2.0, 3.8, 2.2,
    "Logic: Spring Boot",
    [
        "ebank-api本体",
        "業務ログ/監査ログ/セキュリティログ",
        "Path: /opt/ebank-api/app/logs",
    ],
    fill_color=COLOR_APP,
)

add_card(
    slide, 9.1, 2.0, 3.5, 2.2,
    "DB: PostgreSQL 16",
    [
        "bank_appデータベース",
        "DBサーバログを記録",
        "Path: /var/lib/postgresql/16/main/log",
    ],
    fill_color=COLOR_DB,
)

add_line(slide, 4.5, 3.1, 4.9, 3.1)
add_line(slide, 8.7, 3.1, 9.1, 3.1)
add_label(slide, 4.56, 2.8, "HTTP/API")
add_label(slide, 8.73, 2.8, "JDBC")

add_card(
    slide, 1.6, 4.75, 10.2, 1.65,
    "ログ保存先(要点)",
    [
        "Apache: access.log / error.log    Spring Boot: api.log / audit.log / security.log    PostgreSQL: postgresql-*.log",
    ],
)


# 3. Apache logs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "2. Webサーバログ (Apache)")
add_card(
    slide, 0.8, 2.0, 12.0, 4.6,
    "出力先と内容",
    [
        "/var/log/apache2/access.log",
        "  - リクエスト元IP、時刻、メソッド、URI、ステータス、レスポンスサイズ",
        "/var/log/apache2/error.log",
        "  - Apache自体のエラー、設定不整合、上流接続エラーなど",
        "確認コマンド例: tail -f /var/log/apache2/access.log /var/log/apache2/error.log",
    ],
    fill_color=COLOR_WEB,
)


# 4. Spring Boot logs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "3. ロジックログ (Spring Boot / ebank-api)")
add_card(
    slide, 0.8, 1.95, 12.0, 5.0,
    "/opt/ebank-api/app/logs 配下",
    [
        "api.log",
        "  - アプリ通常ログ (INFO/WARN/ERROR)。requestId付きで出力",
        "audit.log",
        "  - 監査ログ(JSON)。業務操作の成功/失敗、actor、path、status など",
        "security.log",
        "  - セキュリティイベント(JSON)。認証関連イベント等",
        "補足(systemd): /opt/ebank-api/logs/app.out.log と app.err.log に標準出力/標準エラー",
    ],
    fill_color=COLOR_APP,
)


# 5. DB logs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "4. DBログ (PostgreSQL)")
add_card(
    slide, 0.8, 2.0, 12.0, 4.6,
    "出力先と内容",
    [
        "/var/lib/postgresql/16/main/log/",
        "  - 接続/切断ログ",
        "  - SQL実行ログ(設定に応じてDDL、遅延SQL、エラーSQLなど)",
        "  - 監査拡張(pgaudit利用時)の監査イベント",
        "確認コマンド例: ls -ltr /var/lib/postgresql/16/main/log/",
    ],
    fill_color=COLOR_DB,
)


# 6. Correlation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "5. ログを見る順番 (運用向け)")
add_bullets(slide, 0.9, 2.0, 12.0, 4.8, [
    "1) まず Apache access.log で対象リクエストを特定 (時刻・IP・URI・status)",
    "2) 同時刻の api.log / audit.log / security.log でアプリ処理結果を確認",
    "3) 必要に応じて PostgreSQLログでSQLエラーや遅延を確認",
    "4) 3層で時刻を突き合わせて、事象を一本の時系列にする",
], size=20)


# 7. Recommended points
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "6. 追加で言っておくと良い運用ポイント")
add_bullets(slide, 0.9, 1.95, 12.0, 5.2, [
    "サーバ時刻同期(NTP)を全層で有効化する (時系列分析の前提)",
    "ログローテーションと保存期間を明確化する (容量枯渇・証跡欠損を防止)",
    "ログ閲覧権限を最小化する (個人情報・認証情報の露出抑止)",
    "運用手順に『どのログを先に見るか』を明記する (初動短縮)",
    "可能なら集中管理(SIEM等)へ転送して横断検索できるようにする",
], size=20)

# 8. Appendix Apache
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "付録A: Apacheログが出るタイミング")
add_card(
    slide, 0.8, 1.85, 6.0, 4.9,
    "出力タイミング",
    [
        "access.log (blue_combined):",
        "  - HTTPリクエスト完了時に1行出力",
        "  - reqid=%{X-Request-Id}i, rt_us=%D を含む",
        "access_error_only.log:",
        "  - REQUEST_STATUS >= 400 の時だけ出力",
        "error.log:",
        "  - Apache内部エラー、上流接続失敗、設定警告時",
    ],
    fill_color=COLOR_WEB,
)
add_card(
    slide, 7.05, 1.85, 5.3, 4.9,
    "追記した主な設定",
    [
        "ServerTokens Prod / ServerSignature Off / TraceEnable Off",
        "LogFormat ... blue_combined",
        "CustomLog \"logs/access_error_only.log\" ... env=is_error",
        "ErrorLog \"logs/error.log\" / LogLevel warn",
        "remoteip_module:",
        "  - RemoteIPHeader X-Forwarded-For",
        "  - RemoteIPInternalProxy 127.0.0.1, ::1, 10/8, 172.16/12, 192.168.1.59",
    ],
    fill_color=COLOR_BOX,
)

# 9. Appendix Spring Boot
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "付録B: Spring Bootログが出るタイミング (ソース準拠)")
add_card(
    slide, 0.8, 1.85, 12.0, 4.9,
    "/opt/ebank-api/app/logs",
    [
        "api.log:",
        "  - 全リクエスト開始時: RequestIdFilter が requestId を採番・記録",
        "  - 各Controllerの処理開始/主要処理でINFO出力",
        "  - 例外時: GlobalExceptionHandler が 4xx=warn / 5xx=error を出力",
        "audit.log:",
        "  - 業務APIの成功時 audit.success(...) で1イベント出力",
        "  - 業務APIの失敗時 audit.fail(...) で1イベント出力",
        "security.log:",
        "  - JWT期限切れ/不正トークン検知時",
        "  - ログイン失敗バーストや異常イベント検知時(sec.emit)",
    ],
    fill_color=COLOR_APP,
)

# 10. Appendix DB
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header(slide, "付録C: PostgreSQLログが出るタイミング")
add_card(
    slide, 0.8, 1.85, 6.0, 4.9,
    "出力タイミング",
    [
        "接続時/切断時:",
        "  - log_connections=on, log_disconnections=on",
        "DDL実行時:",
        "  - log_statement='ddl'",
        "遅延SQL:",
        "  - log_min_duration_statement=500ms 超で出力",
        "エラーSQL:",
        "  - log_min_error_statement=error",
        "ロック待ち/チェックポイント:",
        "  - log_lock_waits=on, log_checkpoints=on",
    ],
    fill_color=COLOR_DB,
)
add_card(
    slide, 7.05, 1.85, 5.3, 4.9,
    "追記した主な設定",
    [
        "logging_collector=on / log_destination='stderr'",
        "log_directory='log'",
        "log_filename='postgresql-%Y-%m-%d.log'",
        "log_rotation_age=1d / log_rotation_size=100MB",
        "log_line_prefix='%m [%p] user=%u,db=%d,app=%a,client=%r,tx=%x '",
        "log_truncate_on_rotation=on",
        "log_error_verbosity=default",
    ],
    fill_color=COLOR_BOX,
)

prs.save(OUT)
print(f"generated: {OUT.resolve()} (slides={len(prs.slides)})")
